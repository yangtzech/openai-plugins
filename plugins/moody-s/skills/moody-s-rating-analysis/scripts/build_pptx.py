#!/usr/bin/env python3
"""Build a Moody's Rating Analysis Deck by populating the official Moody's Corp
Template PPTX (assets/Moody_Corp_Template.pptx).

Rather than building slides from scratch, this script opens the branded
template, removes its demo slides, and adds slides that use the template's
built-in layouts (Cover, Agenda, Divider, 1/2/3 Column, Back Cover,
Disclaimer). All fonts, colors, and decorative elements come from the
template's master slides — we only fill placeholders and add charts/tables.

Usage:
    python build_pptx.py <payload.json> <output.pptx>
"""
from __future__ import annotations

import json
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (XL_CHART_TYPE, XL_LABEL_POSITION,
                             XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION)
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


TEMPLATE_PATH = (Path(__file__).resolve().parent.parent
                 / "assets" / "Moody_Corp_Template.pptx")


# ---- Layout indices (from Moody's Corp Template_2026) ----------------------

L_COVER = 0            # Cover 1
L_AGENDA = 10          # Agenda 4 (3 object columns)
L_DIVIDER = 11         # Divider 1 - Short title
L_ONE_COL = 24         # 1 Column - Subhead
L_TWO_COL = 25         # 2 Column, Equal - Subhead
L_THREE_COL = 27       # 3 Column - Subhead
L_FOUR_COL = 28        # 4 Column - Subhead
L_TITLE_ONLY = 44      # Title Only
L_BACK_COVER = 57      # Back Cover 1 (Thank you)
L_DISCLAIMER = 63      # Disclaimer

# Brand palette used for chart accents (matches template theme colors).
NAVY = RGBColor(0x0A, 0x12, 0x64)
BRIGHT_BLUE = RGBColor(0x00, 0x5E, 0xFF)
MID_BLUE = RGBColor(0x35, 0x39, 0x7E)
PURPLE_BLUE = RGBColor(0x76, 0x77, 0xA7)
PALE_PURPLE = RGBColor(0xA2, 0xA2, 0xC4)
VERY_PALE = RGBColor(0xCF, 0xD0, 0xE1)
PINK = RGBColor(0xD9, 0x01, 0x7A)
LIGHT_GRAY = RGBColor(0xD7, 0xD8, 0xD7)
TEAL = RGBColor(0x4F, 0xB3, 0xAE)
GOLD = RGBColor(0xC4, 0xA5, 0x1A)
GRID_LINE = RGBColor(0xCF, 0xD0, 0xE1)

PIE_PALETTE = [BRIGHT_BLUE, TEAL, GOLD, MID_BLUE, PURPLE_BLUE, PINK,
               PALE_PURPLE, VERY_PALE]
OUTLOOK_PALETTE = {
    "stable": LIGHT_GRAY,
    "positive": TEAL,
    "negative": GOLD,
    "under review": BRIGHT_BLUE,
}

TEXT_DARK = RGBColor(0x22, 0x22, 0x22)
TEXT_MUTED = RGBColor(0x66, 0x6B, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_FILL = RGBColor(0xE9, 0xEA, 0xEE)

# Full Moody's rating scale: numeric notch → label (Aaa=21 … C=1)
RATING_SCALE = {
    21: "Aaa",
    20: "Aa1", 19: "Aa2", 18: "Aa3",
    17: "A1",  16: "A2",  15: "A3",
    14: "Baa1", 13: "Baa2", 12: "Baa3",
    11: "Ba1",  10: "Ba2",   9: "Ba3",
     8: "B1",    7: "B2",    6: "B3",
     5: "Caa1",  4: "Caa2",  3: "Caa3",
     2: "Ca",    1: "C",
}

def notch_to_label(n):
    return RATING_SCALE.get(int(round(n)), str(int(round(n))))


# ---- Utility helpers -------------------------------------------------------

def remove_all_slides(prs):
    """Keep the template's masters/layouts; drop all existing slides."""
    id_lst = prs.slides._sldIdLst
    for sld in list(id_lst):
        rid = sld.attrib[qn("r:id")]
        prs.part.drop_rel(rid)
        id_lst.remove(sld)


_TAG_RE = __import__("re").compile(r"<[^>]+>")


def html_to_text(s):
    """Strip HTML tags and unescape a small set of entities. Paragraphs and
    list items are converted to blank-line separated / bulleted plain text."""
    if s is None:
        return ""
    s = str(s)
    s = (s.replace("</p>", "\n\n")
          .replace("<br>", "\n")
          .replace("<br/>", "\n")
          .replace("<br />", "\n")
          .replace("<li>", "• ")
          .replace("</li>", "\n"))
    s = _TAG_RE.sub("", s)
    s = (s.replace("&amp;", "&")
          .replace("&nbsp;", " ")
          .replace("&rarr;", "→")
          .replace("&#39;", "'")
          .replace("&quot;", '"')
          .replace("&lt;", "<")
          .replace("&gt;", ">"))
    return "\n".join(line.rstrip() for line in s.splitlines()).strip()


def get_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def del_ph(slide, idx):
    ph = get_ph(slide, idx)
    if ph is not None:
        ph._element.getparent().remove(ph._element)


def clear_empty_placeholders(slide):
    """Remove any placeholder that was left empty so PowerPoint doesn't show
    'Click to add text' prompts."""
    for ph in list(slide.placeholders):
        try:
            if ph.has_text_frame and not ph.text_frame.text.strip():
                ph._element.getparent().remove(ph._element)
        except Exception:
            pass


def ph_bounds(slide, idx):
    ph = get_ph(slide, idx)
    if ph is None:
        return None
    return (ph.left, ph.top, ph.width, ph.height)


def set_text(slide, idx, text):
    ph = get_ph(slide, idx)
    if ph is None:
        return None
    cleaned = html_to_text(text)
    if not cleaned:
        del_ph(slide, idx)
        return None
    tf = ph.text_frame
    tf.text = cleaned
    if idx == 0:
        try:
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            size_pt = None
            if len(cleaned) > 60:
                size_pt = 18
            elif len(cleaned) > 45:
                size_pt = 22
            elif len(cleaned) > 35:
                size_pt = 26
            if size_pt is not None:
                for p in tf.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(size_pt)
        except Exception:
            pass
    return ph


def title_bottom(slide, default=Inches(1.1), text=None):
    """Return the bottom edge (EMU) of the title placeholder so downstream
    content can be positioned beneath it regardless of template layout.
    If ``text`` is provided, adds a safety buffer when the title is likely
    to wrap based on a conservative chars-per-line estimate."""
    ph = get_ph(slide, 0)
    if ph is None:
        return default
    try:
        bot = ph.top + ph.height
        if text:
            chars_per_line = max(20, int(ph.width / Inches(0.17)))
            extra_lines = max(0, (len(text) - 1) // chars_per_line)
            if extra_lines > 0:
                bot += Inches(0.45) * extra_lines
        return bot
    except Exception:
        return default


def set_bullets(slide, idx, items, *, bullet_char=None, size=None):
    """Populate a placeholder with multiple paragraphs. Formatting comes
    from the layout/master — we only set the text content."""
    ph = get_ph(slide, idx)
    if ph is None:
        return None
    tf = ph.text_frame
    tf.word_wrap = True
    cleaned = [html_to_text(x) for x in (items or []) if x is not None]
    cleaned = [c for c in cleaned if c]
    if not cleaned:
        del_ph(slide, idx)
        return None
    tf.text = cleaned[0]
    for item in cleaned[1:]:
        p = tf.add_paragraph()
        p.text = item
    if size is not None:
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(size)
    return ph


def set_bold_prefix_bullets(slide, idx, items, *, size=None):
    """Items are [{bold, text}, ...] — each becomes one paragraph with the
    prefix in bold followed by the rest of the text."""
    ph = get_ph(slide, idx)
    if ph is None:
        return None
    tf = ph.text_frame
    tf.word_wrap = True
    tf.text = ""
    p0 = tf.paragraphs[0]
    for i, item in enumerate(items or []):
        if isinstance(item, dict):
            prefix = (item.get("bold") or item.get("prefix")
                      or item.get("date") or "")
            rest = (item.get("text") or item.get("summary") or "")
        elif isinstance(item, (tuple, list)) and len(item) >= 2:
            prefix, rest = str(item[0]), str(item[1])
        else:
            prefix, rest = "", str(item)
        p = p0 if i == 0 else tf.add_paragraph()
        # Clear any inherited run, then add prefix + rest as two runs.
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        if prefix:
            r1 = p.add_run()
            r1.text = prefix + (" – " if rest and not prefix.endswith(("–", ":")) else " ")
            r1.font.bold = True
            if size is not None:
                r1.font.size = Pt(size)
        if rest:
            r2 = p.add_run()
            r2.text = rest
            if size is not None:
                r2.font.size = Pt(size)


# ---- Chart + table helpers -------------------------------------------------

def _style_axes(chart, *, size=8):
    for axis in (chart.category_axis, chart.value_axis):
        try:
            axis.tick_labels.font.name = "Arial"
            axis.tick_labels.font.size = Pt(size)
            axis.tick_labels.font.color.rgb = TEXT_MUTED
        except Exception:
            pass
        try:
            if axis.has_major_gridlines:
                axis.major_gridlines.format.line.color.rgb = GRID_LINE
                axis.major_gridlines.format.line.width = Pt(0.5)
        except Exception:
            pass


def _style_title(chart, title, *, size=13):
    if not title:
        chart.has_title = False
        return
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    for p in chart.chart_title.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = TEXT_DARK


def add_bar_chart(slide, left, top, width, height, title, categories, values,
                  *, color=BRIGHT_BLUE, y_format=None, show_value=True):
    cd = CategoryChartData()
    cd.categories = [str(c) for c in categories]
    cd.add_series(title or "Series",
                  [float(v) if v is not None else 0.0 for v in values])
    title_h = Inches(0.3) if title else 0
    chart_top = top + title_h
    chart_h = height - title_h
    if title:
        add_textbox(slide, left, top, width, title_h, title,
                    size=12, bold=True, color=TEXT_DARK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                left, chart_top, width, chart_h, cd)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = show_value
    if show_value:
        dl = plot.data_labels
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
        dl.font.name = "Arial"
        dl.font.size = Pt(9)
        dl.font.bold = True
        dl.font.color.rgb = TEXT_DARK
        if y_format:
            dl.number_format = y_format
            dl.number_format_is_linked = False
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.format.line.fill.background()
    _style_axes(chart)
    try:
        chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    except Exception:
        pass
    return chart


def add_line_chart(slide, left, top, width, height, title, categories, values,
                   *, color=BRIGHT_BLUE, rating_scale=False):
    cd = CategoryChartData()
    cd.categories = [str(c) for c in categories]
    cd.add_series(title or "Series",
                  [float(v) if v is not None else None for v in values])

    if rating_scale:
        # Reserve left margin for custom rating label textboxes
        label_w = Inches(0.55)
        chart_left  = left + label_w
        chart_width = width - label_w
    else:
        chart_left  = left
        chart_width = width

    title_h = Inches(0.3) if title else 0
    if title:
        add_textbox(slide, left, top, width, title_h, title,
                    size=12, bold=True, color=TEXT_DARK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE,
                                chart_left, top + title_h,
                                chart_width, height - title_h, cd)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = False
    s = chart.series[0]
    s.format.line.color.rgb = color
    s.format.line.width = Pt(2)
    try:
        s.marker.size = 6
        s.marker.format.fill.solid()
        s.marker.format.fill.fore_color.rgb = color
        s.marker.format.line.color.rgb = color
    except Exception:
        pass
    _style_axes(chart)

    if rating_scale:
        try:
            clean = [float(v) for v in values if v is not None]
            if clean:
                # ── Smart scale: anchors + contextual window ──────────────
                current = int(round(max(clean)))
                ANCHORS = {21, 18, 15, 12, 9, 6, 3}  # Aaa,Aa3,A3,Baa3,Ba3,B3,Caa3
                mn = max(1,  current - 5)
                mx = min(21, current + 5)
                context = set(range(max(1, current - 2), min(21, current + 2) + 1))
                visible = (ANCHORS | context) & set(range(mn, mx + 1))

                val_ax = chart.value_axis
                val_ax.minimum_scale = mn
                val_ax.maximum_scale = mx
                val_ax.major_unit    = 1

                # Hide the numeric axis labels — PowerPoint ignores conditional
                # numFmt on value axes, so we use a textbox overlay instead.
                try:
                    val_ax.tick_labels.font.size = Pt(1)
                    val_ax.tick_labels.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                except Exception:
                    pass
                from lxml import etree
                ns      = "http://schemas.openxmlformats.org/drawingml/2006/chart"
                fmt_tag = f"{{{ns}}}numFmt"
                existing = val_ax._element.find(fmt_tag)
                if existing is not None:
                    existing.attrib["formatCode"]  = '""'
                    existing.attrib["sourceLinked"] = "0"
                else:
                    el = etree.SubElement(val_ax._element, fmt_tag,
                                          attrib={"formatCode": '""',
                                                  "sourceLinked": "0"})
                    val_ax._element.remove(el)
                    val_ax._element.insert(1, el)

                # ── Textbox overlay: place rating labels at correct Y positions ──
                chart_top_emu   = top + title_h
                chart_h_emu     = height - title_h
                # Approximate PowerPoint plot area vertical bounds
                plot_top    = chart_top_emu + int(chart_h_emu * 0.08)
                plot_bottom = chart_top_emu + int(chart_h_emu * 0.88)
                plot_height = plot_bottom - plot_top
                lbl_h = Inches(0.18)
                lbl_w = label_w - Emu(36000)

                for notch in range(mn, mx + 1):
                    if notch not in visible:
                        continue
                    frac     = (notch - mn) / (mx - mn)
                    y_centre = plot_bottom - int(plot_height * frac)
                    y        = y_centre - lbl_h // 2
                    label    = notch_to_label(notch)
                    is_anchor  = notch in ANCHORS
                    is_current = notch == current
                    tb = slide.shapes.add_textbox(left, y, lbl_w, lbl_h)
                    tf = tb.text_frame
                    tf.word_wrap = False
                    tf.margin_left = tf.margin_right = 0
                    tf.margin_top  = tf.margin_bottom = 0
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.RIGHT
                    r = p.add_run()
                    r.text = label
                    r.font.name  = "Arial"
                    r.font.size  = Pt(9 if is_anchor else 8)
                    r.font.bold  = is_anchor or is_current
                    r.font.color.rgb = (NAVY if is_current
                                        else TEXT_DARK if is_anchor
                                        else TEXT_MUTED)
        except Exception:
            pass
    return chart


def add_pie_chart(slide, left, top, width, height, title, labels, values,
                  *, colors=None, show_percent=True, use_legend=False,
                  legend_position=None):
    cd = CategoryChartData()
    cd.categories = list(labels)
    cd.add_series(title or "Series", [float(v or 0) for v in values])
    title_h = Inches(0.35) if title else 0
    if title:
        add_textbox(slide, left, top, width, title_h, title,
                    size=12, bold=True, color=TEXT_DARK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.PIE,
                                left, top + title_h, width,
                                height - title_h, cd)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = bool(use_legend)
    if use_legend:
        try:
            chart.legend.position = legend_position or XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.name = "Arial"
            chart.legend.font.size = Pt(9)
            chart.legend.font.color.rgb = TEXT_DARK
        except Exception:
            pass
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_percentage = show_percent
    dl.show_value = not show_percent
    dl.show_category_name = not use_legend
    dl.font.name = "Arial"
    dl.font.size = Pt(9)
    dl.font.bold = True
    dl.font.color.rgb = TEXT_DARK
    dl.position = (XL_LABEL_POSITION.BEST_FIT if use_legend
                   else XL_LABEL_POSITION.OUTSIDE_END)
    palette = list(colors) if colors else PIE_PALETTE
    for i, pt in enumerate(chart.series[0].points):
        c = palette[i % len(palette)]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = c
        pt.format.line.color.rgb = WHITE
        pt.format.line.width = Pt(1)
    return chart


def add_scatter_chart(slide, left, top, width, height, title, points, *,
                      x_label="", y_label=""):
    clean = []
    for p in points or []:
        x = p.get("x")
        y = p.get("y")
        if x is None or y is None:
            continue
        try:
            clean.append((p.get("company", "") or p.get("name", ""),
                          float(x), float(y)))
        except (TypeError, ValueError):
            continue
    if not clean:
        return None, []
    cd = XyChartData()
    s = cd.add_series(title or "Series")
    for _, x, y in clean:
        s.add_data_point(x, y)
    title_h = Inches(0.3) if title else 0
    if title:
        add_textbox(slide, left, top, width, title_h, title,
                    size=12, bold=True, color=TEXT_DARK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER,
                                left, top + title_h, width,
                                height - title_h, cd)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = False
    sr = chart.series[0]
    sr.marker.size = 9
    sr.marker.format.fill.solid()
    sr.marker.format.fill.fore_color.rgb = BRIGHT_BLUE
    sr.marker.format.line.color.rgb = BRIGHT_BLUE
    try:
        sr.format.line.fill.background()
    except Exception:
        pass
    for ax, lab in ((chart.category_axis, x_label),
                     (chart.value_axis, y_label)):
        try:
            ax.has_title = bool(lab)
            if lab:
                ax.axis_title.text_frame.text = lab
                for p in ax.axis_title.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = "Arial"
                        r.font.size = Pt(9)
                        r.font.color.rgb = TEXT_MUTED
            ax.tick_labels.font.name = "Arial"
            ax.tick_labels.font.size = Pt(8)
            ax.tick_labels.font.color.rgb = TEXT_MUTED
        except Exception:
            pass
        try:
            if ax.has_major_gridlines:
                ax.major_gridlines.format.line.color.rgb = GRID_LINE
                ax.major_gridlines.format.line.width = Pt(0.5)
        except Exception:
            pass
    return chart, clean


def add_scatter_labels(slide, left, top, width, height, clean,
                       *, title_offset=Inches(0.3)):
    """Place small text labels above each scatter point. The bounds passed in
    are the original slot (same as add_scatter_chart); we apply the same title
    offset the chart helper uses so labels line up with the plot."""
    if not clean:
        return
    chart_top = top + title_offset
    chart_h = height - title_offset
    xs = [c[1] for c in clean]
    ys = [c[2] for c in clean]
    xmin, xmax = min(xs), max(xs)
    ymin, ymax = min(ys), max(ys)
    if xmax == xmin:
        xmax = xmin + 1
    if ymax == ymin:
        ymax = ymin + 1
    plot_left = left + int(width * 0.14)
    plot_top = chart_top + int(chart_h * 0.08)
    plot_w = int(width * 0.80)
    plot_h = int(chart_h * 0.74)
    placed = []
    label_w = Inches(0.95)
    label_h = Inches(0.18)
    for name, x, y in clean:
        if not name:
            continue
        fx = (x - xmin) / (xmax - xmin)
        fy = 1 - (y - ymin) / (ymax - ymin)
        px = plot_left + int(plot_w * fx)
        py = plot_top + int(plot_h * fy)
        lx = px - label_w // 2
        ly = py - label_h - Inches(0.08)
        for _ in range(6):
            collision = any(
                abs(lx - ox) < label_w * 0.75 and abs(ly - oy) < label_h
                for ox, oy in placed
            )
            if not collision:
                break
            ly -= label_h
        placed.append((lx, ly))
        tb = slide.shapes.add_textbox(lx, ly, label_w, label_h)
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = name
        r.font.name = "Arial"
        r.font.size = Pt(8)
        r.font.bold = True
        r.font.color.rgb = TEXT_DARK


def add_textbox(slide, left, top, width, height, text, *, size=10, bold=False,
                color=TEXT_DARK, font="Arial", align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text or "")
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_table(slide, left, top, width, height, headers, rows, *,
              header_fill=HEADER_FILL, first_col_bold=False,
              header_size=10, body_size=9, col_widths=None):
    nrows = len(rows) + 1
    ncols = len(headers)
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    table = shape.table
    if col_widths:
        total = float(sum(col_widths))
        for i, w in enumerate(col_widths):
            table.columns[i].width = int(width * (w / total))
    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = header_fill
        tf = c.text_frame
        tf.text = str(h or "")
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.name = "Arial"
                r.font.size = Pt(header_size)
                r.font.bold = True
                r.font.color.rgb = TEXT_DARK
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE
            tf = c.text_frame
            tf.text = str(val) if val is not None else ""
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.name = "Arial"
                    r.font.size = Pt(body_size)
                    r.font.bold = first_col_bold and j == 0
                    r.font.color.rgb = TEXT_DARK
    return table


# ---- HTML parsing for news/trends sections ---------------------------------

def _parse_html_groups(html: str):
    """Parse '<strong>Group</strong><ul><li>item</li>...</ul>' chunks."""
    if not html:
        return []
    out = []
    for chunk in html.split("<strong")[1:]:
        _, _, rest = chunk.partition(">")
        name, _, tail = rest.partition("</strong>")
        items = []
        if "<ul>" in tail and "</ul>" in tail:
            ul = tail.split("<ul>", 1)[1].split("</ul>", 1)[0]
            for li in ul.split("<li>")[1:]:
                t = li.split("</li>", 1)[0].strip()
                t = (t.replace("&amp;", "&")
                      .replace("&nbsp;", " ")
                      .replace("&rarr;", "→")
                      .replace("&#39;", "'"))
                if t:
                    items.append(t)
        name = name.replace("&amp;", "&").strip()
        out.append({"name": name, "items": items})
    return out


def _normalize_categorized(section):
    """Return list of {name, items} from string/list/dict input."""
    if isinstance(section, str):
        return _parse_html_groups(section)
    if isinstance(section, list):
        out = []
        for g in section:
            if isinstance(g, dict):
                out.append({
                    "name": g.get("name") or g.get("category") or "",
                    "items": g.get("items") or [],
                })
        return out
    if isinstance(section, dict):
        if section.get("groups"):
            return _normalize_categorized(section["groups"])
        if section.get("html"):
            return _parse_html_groups(section["html"])
    return []


# ---- Slide builders --------------------------------------------------------

def build_cover(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_COVER])
    company = data.get("target_company", "")
    set_text(s, 0, company)
    set_text(s, 100, "Rating Analysis DECK")
    set_text(s, 101, data.get("sector", ""))
    set_text(s, 102, f"Report date: {data.get('report_date', '')}")
    set_text(s, 99, "MOODY'S ANALYTICS")
    return s


def build_agenda(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_AGENDA])
    agenda = [
        ("01", "Sector Analysis",
         "Macro context, sector outlook and recent rating actions"),
        ("02", "Company Credit Overview",
         "Financial profile, SWOT, strategy and news"),
        ("03", "Company Positioning vs. Peers",
         "Benchmarking on leverage, profitability, scorecard and ESG"),
    ]
    for i, (num, title, desc) in enumerate(agenda):
        ph = get_ph(s, 18 + i)
        if ph is None:
            continue
        tf = ph.text_frame
        tf.word_wrap = True
        tf.text = ""
        p0 = tf.paragraphs[0]
        r = p0.add_run()
        r.text = num
        r.font.name = "Times New Roman"
        r.font.size = Pt(64)
        r.font.bold = True
        r.font.color.rgb = BRIGHT_BLUE
        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = title
        r1.font.name = "Arial"
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = TEXT_DARK
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = desc
        r2.font.name = "Arial"
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_MUTED
    return s


def build_divider(prs, title):
    s = prs.slides.add_slide(prs.slide_layouts[L_DIVIDER])
    set_text(s, 0, title.upper())
    return s


def build_sector_overview(prs, data):
    """3-column layout: Overview | Watchlist | Takeaways."""
    s = prs.slides.add_slide(prs.slide_layouts[L_THREE_COL])
    sec = data.get("sections", {}).get("sector_overview", {}) or {}
    sector = data.get("sector", "")
    set_text(s, 0, f"{sector} – Sector overview")
    set_text(s, 21, "")
    set_text(s, 29, "SECTOR OVERVIEW")
    set_text(s, 30, "WHAT TO LOOK OUT FOR?")
    set_text(s, 31, "KEY TAKEAWAYS")
    set_text(s, 19, "Source: Moodys.com")
    # Convert object placeholders into text content.
    set_bullets(s, 23, sec.get("overview_bullets", []), size=11)
    set_bullets(s, 24, sec.get("watchlist_bullets", []), size=11)
    set_bullets(s, 25, sec.get("takeaway_bullets", []), size=11)
    return s


def build_moodys_view(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_TWO_COL])
    mv = data.get("sections", {}).get("moodys_view", {}) or {}
    sector = data.get("sector", "")
    title_text = f"{sector} – Moody's view"
    set_text(s, 0, title_text)
    set_text(s, 19, "Source: Moodys.com")
    del_ph(s, 21)
    del_ph(s, 29)
    del_ph(s, 30)
    bnd_l = ph_bounds(s, 27)
    bnd_r = ph_bounds(s, 28)
    del_ph(s, 27)
    del_ph(s, 28)

    summary = html_to_text(mv.get("outlook_summary"))
    pos = html_to_text(mv.get("company_positioning"))
    dist = mv.get("outlook_distribution") or []
    dist2 = mv.get("ratings_distribution") or []

    if not bnd_l or not bnd_r:
        title_bot = title_bottom(s, default=Inches(1.1), text=title_text)
        fallback_top = max(title_bot + Inches(0.2), Inches(1.3))
        fallback_h = Inches(5.5)
        if not bnd_l:
            bnd_l = (Inches(0.6), fallback_top, Inches(6.0), fallback_h)
        if not bnd_r:
            bnd_r = (Inches(6.9), fallback_top, Inches(6.0), fallback_h)

    if summary or pos:
        l, t, w, h = bnd_l
        tb = s.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.text = ""
        lines = [x for x in [summary, pos] if x]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = 1.2
            p.space_after = Pt(6)
            r = p.add_run()
            r.text = line
            r.font.name = "Arial"
            r.font.size = Pt(11)
            r.font.color.rgb = TEXT_DARK

    pie_bounds = bar_bounds = None
    if dist and dist2:
        l, t, w, h = bnd_r
        half = int(h / 2)
        pie_bounds = (l, t, w, half)
        bar_bounds = (l, t + half, w, h - half)
    elif dist:
        pie_bounds = bnd_r
    elif dist2:
        bar_bounds = bnd_r

    if pie_bounds and dist:
        l, t, w, h = pie_bounds
        labels, values, colors = [], [], []
        for d in dist:
            cat = str(d.get("category", ""))
            pct = d.get("percentage", d.get("count", 0))
            count = d.get("count", pct or 0)
            try:
                if count is None or float(count) <= 0:
                    continue
            except (TypeError, ValueError):
                continue
            labels.append(cat)
            values.append(count)
            colors.append(OUTLOOK_PALETTE.get(cat.lower(), BRIGHT_BLUE))
        if values:
            add_pie_chart(s, l, t, w, h,
                          "Moody's Outlooks Distribution",
                          labels, values, colors=colors, use_legend=True)

    if bar_bounds and dist2:
        l, t, w, h = bar_bounds
        add_bar_chart(s, l, t, w, h,
                      "Moody's Ratings Distribution",
                      [str(d.get("rating", "")) for d in dist2],
                      [d.get("count", 0) for d in dist2],
                      color=BRIGHT_BLUE)
    return s


def build_macro_outlook(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_ONE_COL])
    set_text(s, 0, "Global Macro Outlook")
    set_text(s, 18, "MOODY'S REAL GDP GROWTH")
    macro = data.get("sections", {}).get("macro_outlook", {}) or {}
    set_text(s, 29, "")
    set_text(s, 19, "Source: Moodys.com")
    # Replace the content placeholder with bullets above a GDP table.
    bnd = ph_bounds(s, 22)
    del_ph(s, 22)
    if not bnd:
        return s
    l, t, w, h = bnd
    new_t = Inches(1.65)
    h = h + (t - new_t)
    t = new_t
    bullet_h = Inches(1.2)
    add_textbox_bullets(s, l, t, w, bullet_h,
                        macro.get("gdp_commentary", []), size=10)
    tbl = macro.get("gdp_table", {}) or {}
    years = tbl.get("year_columns", []) or []
    rows = tbl.get("rows", []) or []
    if years and rows:
        headers = ["Country"] + [f"{y} GDP Growth" for y in years]
        body = [[r.get("country", "")] + list(r.get("values", []))
                for r in rows]
        add_table(s, l, t + bullet_h, w, h - bullet_h, headers, body,
                  col_widths=[2.0] + [1.2] * len(years),
                  header_size=10, body_size=9)
    return s


def _fit_bullet_font(width_emu, height_emu, items, size, bullet="•"):
    """Pick a font size (<= ``size``, >= 8) that makes the bullet list fit
    vertically within ``height_emu``. Uses a conservative chars-per-line
    estimate for Arial and iterates 1pt at a time."""
    # Coerce bare string so we don't iterate characters
    if isinstance(items, str):
        items = [items] if items.strip() else []
    try:
        width_pt = (float(width_emu) / Inches(1)) * 72.0
        avail_pt = (float(height_emu) / Inches(1)) * 72.0
    except Exception:
        return size
    if width_pt <= 0 or avail_pt <= 0 or not items:
        return size
    pad_pt = 14.4
    for sz in range(int(size), 7, -1):
        avg_char_pt = sz * 0.52
        effective_pt = max(1.0, width_pt - pad_pt)
        cpl = max(10, int(effective_pt / avg_char_pt))
        total_lines = 0
        for item in items:
            text = f"{bullet} {item}"
            total_lines += max(1, -(-len(text) // cpl))
        line_spacing = 1.2 if sz <= 9 else 1.25
        space_after = 2 if sz <= 9 else 4
        used_pt = total_lines * sz * line_spacing + len(items) * space_after
        if used_pt <= avail_pt * 0.95:
            return sz
    return 8


def add_textbox_bullets(slide, left, top, width, height, items, *,
                        size=11, bullet="•"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = ""
    # Guard: a bare string must be wrapped in a list; iterating a string
    # directly produces one character per bullet (the "• C \n • o \n …" bug).
    if isinstance(items, str):
        items = [items] if items.strip() else []
    items = [str(x) for x in (items or []) if x is not None and str(x).strip()]
    if not items:
        return tb
    fit_size = _fit_bullet_font(width, height, items, size, bullet=bullet)
    line_spacing = 1.2 if fit_size <= 9 else 1.25
    space_after_pt = 2 if fit_size <= 9 else 4
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after_pt)
        r0 = p.add_run()
        r0.text = f"{bullet} "
        r0.font.name = "Arial"
        r0.font.size = Pt(fit_size)
        r0.font.color.rgb = TEXT_DARK
        r1 = p.add_run()
        r1.text = item
        r1.font.name = "Arial"
        r1.font.size = Pt(fit_size)
        r1.font.color.rgb = TEXT_DARK
    return tb


def build_rating_actions(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_ONE_COL])
    set_text(s, 0, "Moody's Rating Actions – YTD")
    set_text(s, 18, "")
    set_text(s, 29, "")
    set_text(s, 19, "Source: Moodys.com")
    actions = data.get("sections", {}).get("rating_actions_ytd", []) or []
    items = []
    for a in actions[:10]:
        date = a.get("date", "")
        company = a.get("company", "")
        summary = a.get("summary", "")
        prefix = f"{date} – {company}" if company else date
        items.append({"bold": prefix, "text": summary})
    set_bold_prefix_bullets(s, 22, items, size=11)
    return s


def build_financial_analysis(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_TWO_COL])
    company = data.get("target_company", "")
    set_text(s, 0, f"{company} financial analysis")
    set_text(s, 21, "")
    set_text(s, 29, "")
    set_text(s, 30, "")
    set_text(s, 19, "Source: Moodys.com, Company filings")
    fa = data.get("sections", {}).get("financial_analysis", {}) or {}
    # Left column (27): bullets.
    bnd_l = ph_bounds(s, 27)
    bnd_r = ph_bounds(s, 28)
    del_ph(s, 27)
    del_ph(s, 28)
    new_t = Inches(1.15)
    if bnd_l:
        l, t, w, h = bnd_l
        h = h + (t - new_t)
        t = new_t
        _fa_commentary = fa.get("commentary", [])
        if isinstance(_fa_commentary, str):
            _fa_commentary = [_fa_commentary] if _fa_commentary.strip() else []
        add_textbox_bullets(s, l, t, w, h, _fa_commentary, size=10)
    if bnd_r:
        l, t, w, h = bnd_r
        h = h + (t - new_t)
        t = new_t
        history = fa.get("rating_history", []) or []
        # ── Sort oldest → newest so chart x-axis reads left-to-right chronologically ──
        _MONTH_MAP = {
            "jan":1,"feb":2,"mar":3,"apr":4,"may":5,"jun":6,
            "jul":7,"aug":8,"sep":9,"oct":10,"nov":11,"dec":12
        }
        def _date_key(entry: dict) -> tuple:
            """Return (year, month) tuple for sorting; fall back to (0,0) if unparseable."""
            raw = entry.get("date", "")
            m = __import__("re").search(r"(\d{4})", raw)
            year = int(m.group(1)) if m else 0
            mn = __import__("re").search(r"(jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)", raw.lower())
            month = _MONTH_MAP.get(mn.group(1), 0) if mn else 0
            return (year, month)
        history = sorted(history, key=_date_key)
        # Rebuild chart_data in the same sorted order (parallel arrays must stay aligned)
        _notch = {"Aaa":21,"Aa1":20,"Aa2":19,"Aa3":18,"A1":17,"A2":16,"A3":15,
                  "Baa1":14,"Baa2":13,"Baa3":12,"Ba1":11,"Ba2":10,"Ba3":9,
                  "B1":8,"B2":7,"B3":6,"Caa1":5,"Caa2":4,"Caa3":3,"Ca":2,"C":1}
        chart_data = [_notch.get(entry.get("rating",""), 0) for entry in history] if history else (fa.get("rating_chart_data", []) or [])
        chart_h = int(h * 0.5)
        if chart_data:
            cats = [entry.get("date", f"P{i+1}") for i, entry in enumerate(history)]
            if len(cats) < len(chart_data):
                cats += [f"P{i+1}" for i in range(len(cats), len(chart_data))]
            elif len(cats) > len(chart_data):
                cats = cats[:len(chart_data)]
            # Single data point: pad with today so the line has direction
            # and reflects that the rating is still current as of report date
            if len(chart_data) == 1:
                import datetime
                today_label = datetime.date.today().strftime("%b %Y")
                cats        = cats + [today_label]
                chart_data  = chart_data + [chart_data[0]]
            add_line_chart(s, l, t, w, chart_h,
                           f"Moody's Rating History for {company}",
                           cats, chart_data, color=BRIGHT_BLUE,
                           rating_scale=True)
        rationale_top = t + chart_h + Inches(0.15)
        add_textbox(s, l, rationale_top, w, Inches(0.3),
                    "Moody's rating rationale", size=11, bold=True,
                    color=TEXT_DARK)
        items = []
        for h2 in history:
            items.append({
                "bold": h2.get("date", ""),
                "text": h2.get("reason", h2.get("rating", "")),
            })
        _add_bold_prefix_textbox(s, l, rationale_top + Inches(0.35), w,
                         Inches(2.4), items, size=10)
    return s


def _add_bold_prefix_textbox(slide, left, top, width, height, items, *,
                             size=10, bullet="•"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = ""
    for i, item in enumerate(items or []):
        if isinstance(item, dict):
            prefix = item.get("bold", "")
            rest = item.get("text", "")
        else:
            prefix, rest = "", str(item)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.2
        p.space_after = Pt(3)
        r0 = p.add_run()
        r0.text = f"{bullet} "
        r0.font.name = "Arial"
        r0.font.size = Pt(size)
        r0.font.color.rgb = TEXT_DARK
        if prefix:
            r1 = p.add_run()
            r1.text = f"{prefix}: "
            r1.font.name = "Arial"
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = TEXT_DARK
        if rest:
            r2 = p.add_run()
            r2.text = rest
            r2.font.name = "Arial"
            r2.font.size = Pt(size)
            r2.font.color.rgb = TEXT_DARK
    return tb


def build_revenue_distribution(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_ONE_COL])
    company = data.get("target_company", "")
    set_text(s, 0, "Revenue distribution")
    set_text(s, 18, "")
    set_text(s, 29, "")
    set_text(s, 19, "Source: Moodys.com, Company filings")
    rd = data.get("sections", {}).get("revenue_distribution", {}) or {}
    bnd = ph_bounds(s, 22)
    del_ph(s, 22)
    if not bnd:
        return s
    l, t, w, h = bnd
    new_t = Inches(1.15)      # move up from 2.653"
    h = h + (t - new_t)      # recover ~1.5" into h
    t = new_t

    bullet_h = Inches(0.6)    # was 1.2 — only 2 short bullets, 0.6" is enough
    add_textbox_bullets(s, l, t, w, bullet_h,
                        rd.get("commentary", []), size=10)
    geo = rd.get("by_geography", []) or []
    seg = rd.get("by_segment", []) or []
    charts_top = t + bullet_h   # shifts up automatically
    charts_h = h - bullet_h     # charts get all the recovered space
    col_w = int((w - Inches(0.4)) / 2)
    def _clean_slices(items):
        out = []
        for it in items or []:
            name = str(it.get("name", "")).strip()
            try:
                pct = float(it.get("percentage", 0) or 0)
            except (TypeError, ValueError):
                pct = 0.0
            if not name or pct <= 0:
                continue
            out.append((name, pct))
        return out

    geo_clean = _clean_slices(geo)
    seg_clean = _clean_slices(seg)

    if geo_clean:
        add_pie_chart(s, l, charts_top, col_w, charts_h,
                      f"{company} Revenue by region",
                      [n for n, _ in geo_clean],
                      [p for _, p in geo_clean],
                      use_legend=True)
    if seg_clean:
        add_pie_chart(s, l + col_w + Inches(0.4), charts_top, col_w,
                      charts_h, f"{company} Revenue by segment",
                      [n for n, _ in seg_clean],
                      [p for _, p in seg_clean],
                      use_legend=True)
    return s


def build_swot(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_ONE_COL])
    set_text(s, 0, "SWOT Analysis")
    set_text(s, 18, "")
    set_text(s, 29, "")
    set_text(s, 19, "Source: Moodys.com, Company filings")
    swot = data.get("sections", {}).get("swot", {}) or {}
    bnd = ph_bounds(s, 22)
    del_ph(s, 22)
    if not bnd:
        return s
    l, t, w, h = bnd
    new_t = Inches(1.15)       
    h = h + (t - new_t)    
    t = new_t
    col_w = int((w - Inches(0.4)) / 2)
    row_h = int((h - Inches(0.2)) / 2)
    quads = [
        ("Strengths", swot.get("strengths", [])),
        ("Weaknesses", swot.get("weaknesses", [])),
        ("Opportunities", swot.get("opportunities", [])),
        ("Threats", swot.get("threats", [])),
    ]
    for i, (label, items) in enumerate(quads):
        col = i % 2
        row = i // 2
        x = l + col * (col_w + Inches(0.4))
        y = t + row * (row_h + Inches(0.2))
        add_textbox(s, x, y, col_w, Inches(0.3), label,
                    size=11, bold=True, color=TEXT_DARK)
        add_textbox_bullets(s, x, y + Inches(0.3), col_w,
                            row_h - Inches(0.3), items, size=10)
    return s


def build_key_metrics(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_ONE_COL])
    company = data.get("target_company", "")
    set_text(s, 0, f"Key financial metrics - {company}")
    set_text(s, 18, "")
    set_text(s, 29, "")
    set_text(s, 19, "Source: Moodys.com, Company filings")
    km = data.get("sections", {}).get("key_metrics", {}) or {}
    periods = km.get("periods", []) or []
    charts = [
        ("Revenue", km.get("revenue", []), "#,##0"),
        ("EBIT Margin", km.get("ebit_margin", []), "0.0\"%\""),
        ("Debt / EBITDA", km.get("debt_ebitda", []), "0.0\"x\""),
        ("RCF / Net Debt", km.get("rcf_net_debt", []), "0.0\"%\""),
    ]
    bnd = ph_bounds(s, 22)
    del_ph(s, 22)
    if not bnd:
        return s
    l, t, w, h = bnd
    new_t = Inches(1.15)
    h = h + (t - new_t)
    t = new_t

    gap = Inches(0.3)
    col_w = int((w - gap) / 2)
    row_h = int((h - gap) / 2)    # ~1.5" taller rows → charts fill the slide properly
    for i, (title, vals, fmt) in enumerate(charts):
        col = i % 2
        row = i // 2
        x = l + col * (col_w + gap)
        y = t + row * (row_h + gap)
        cats = periods[:len(vals)] if periods else [f"P{k+1}" for k in range(len(vals))]
        if vals:
            add_bar_chart(s, x, y, col_w, row_h, title, cats, vals,
                          color=BRIGHT_BLUE, y_format=fmt)
        else:
            add_textbox(s, x, y, col_w, row_h, f"{title}\n(no data)",
                        size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
                        anchor=MSO_ANCHOR.MIDDLE)
    return s


def build_strategic_updates(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_TWO_COL])
    company = data.get("target_company", "")
    set_text(s, 0, f"{company} strategic updates")
    set_text(s, 21, "")
    set_text(s, 29, "RECENT DEVELOPMENTS")
    set_text(s, 30, "FORWARD LOOKING")
    set_text(s, 19, "Source: Moodys.com, Company filings")
    sup = data.get("sections", {}).get("strategic_updates", {}) or {}
    set_bullets(s, 27, sup.get("recent", []), size=11)
    set_bullets(s, 28, sup.get("forward", []), size=11)
    return s


def build_categorized(prs, data, section_key, title, source):
    s = prs.slides.add_slide(prs.slide_layouts[L_ONE_COL])
    set_text(s, 0, title)
    set_text(s, 18, "")
    set_text(s, 29, "")
    set_text(s, 19, f"Source: {source}")
    groups = _normalize_categorized(data.get("sections", {}).get(section_key))
    bnd = ph_bounds(s, 22)
    del_ph(s, 22)
    if not bnd or not groups:
        return s
    l, t, w, h = bnd
    new_t = Inches(1.15)
    h = h + (t - new_t)
    t = new_t

    per_h = int(h / max(len(groups), 1))   # each group gets a proportionally larger slice
    for i, g in enumerate(groups):
        y = t + i * per_h
        add_textbox(s, l, y, w, Inches(0.3), g.get("name", ""),
                    size=11, bold=True, color=TEXT_DARK)
        items = g.get("items") or []
        bullets = []
        for it in items:
            if isinstance(it, dict):
                prefix = it.get("date") or it.get("when") or ""
                body = it.get("text") or it.get("summary") or ""
                bullets.append({"bold": f"{prefix} →" if prefix else "",
                                "text": body})
            else:
                bullets.append(str(it))
        if any(isinstance(b, dict) for b in bullets):
            _add_bold_prefix_textbox(s, l, y + Inches(0.3), w,
                                     per_h - Inches(0.3),
                                     bullets, size=10)
        else:
            add_textbox_bullets(s, l, y + Inches(0.3), w,
                                per_h - Inches(0.3), bullets, size=10)
    return s


def build_peer_summary(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_ONE_COL])
    set_text(s, 0, "Peer comparison summary")
    set_text(s, 18, "")
    set_text(s, 29, "")
    set_text(s, 19, "Source: Moodys.com, Company filings")
    ps = data.get("sections", {}).get("peer_summary", {}) or {}
    table = ps.get("table", []) or []
    bnd = ph_bounds(s, 22)
    del_ph(s, 22)
    if not bnd:
        return s
    l, t, w, h = bnd
    new_t = Inches(1.15)
    h = h + (t - new_t)
    t = new_t

    bul_h = Inches(0.9)
    _ps_commentary = ps.get("commentary", [])
    if isinstance(_ps_commentary, str):
        _ps_commentary = [_ps_commentary] if _ps_commentary.strip() else []
    add_textbox_bullets(s, l, t, w, bul_h, _ps_commentary, size=10)
    if not table:
        return s
    fields = [
        ("Country", "country"),
        ("Market Capitalization", "market_cap"),
        ("Moody's Ratings", "rating"),
        ("Outlook", "outlook"),
        ("Business Mix", "business_mix"),
    ]
    headers = [""] + [c.get("company", "") for c in table]
    body = []
    for label, key in fields:
        body.append([label] + [str(c.get(key, "") or "") for c in table])
    tbl_top = t + bul_h
    tbl_h = h - bul_h
    add_table(s, l, tbl_top, w, tbl_h, headers, body,
              first_col_bold=True,
              col_widths=[1.6] + [1.0] * len(table),
              header_size=10, body_size=9)
    return s


def build_peer_financials(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_ONE_COL])
    set_text(s, 0, "Detailed peer comparison")
    set_text(s, 18, "")
    set_text(s, 29, "")
    set_text(s, 19, "Source: Moodys.com, Company filings")
    pf = data.get("sections", {}).get("peer_financials", {}) or {}
    columns = pf.get("columns", []) or []
    rows = pf.get("rows", []) or []
    bnd = ph_bounds(s, 22)
    del_ph(s, 22)
    if not bnd or not rows:
        return s
    l, t, w, h = bnd

    # Center the table vertically between title (ends ~1.08") and footer (6.583")
    tbl_h = Inches(3.8)
    avail_top = Inches(1.15)
    avail_bot = Inches(6.50)
    t = avail_top + int((avail_bot - avail_top - tbl_h) / 2)  # center it

    headers = ["Financial Items (in millions)"] + [r.get("company", "") for r in rows]
    body = [["Period"] + [str(r.get("period", "") or "") for r in rows],
            ["Currency"] + [str(r.get("currency", "") or "") for r in rows]]
    for i, label in enumerate(columns):
        body.append(
            [label] + [str((r.get("values") or [""] * len(columns))[i]
                       if i < len(r.get("values", [])) else "")
                       for r in rows]
        )
    add_table(s, l, t, w, tbl_h, headers, body,
              first_col_bold=True,
              col_widths=[2.0] + [1.0] * len(rows),
              header_size=10, body_size=9)
    return s


def _build_peer_bars(prs, data, section_key, title, chart_defs):
    s = prs.slides.add_slide(prs.slide_layouts[L_ONE_COL])
    set_text(s, 0, title)
    set_text(s, 18, "")
    set_text(s, 29, "")
    set_text(s, 19, "Source: Moodys.com, Company filings")
    section = data.get("sections", {}).get(section_key, {}) or {}
    bnd = ph_bounds(s, 22)
    del_ph(s, 22)
    if not bnd:
        return s
    l, t, w, h = bnd

    new_t = Inches(1.15)
    h = h + (t - new_t)
    t = new_t

    commentary = section.get("commentary")
    if isinstance(commentary, str):
        commentary = [commentary]
    bul_h = Inches(0.6)          # was 1.0 — commentary is 1-2 lines, 0.6" is enough
    add_textbox_bullets(s, l, t, w, bul_h, commentary or [], size=10)
    charts_top = t + bul_h       # shifts up automatically
    charts_h = h - bul_h        # charts get ~1.9" more height total
    col_w = int((w - Inches(0.4)) / 2)
    for i, (chart_title, key, fmt) in enumerate(chart_defs):
        x = l + i * (col_w + Inches(0.4))
        entries = section.get(key, []) or []
        clean = [(e.get("company", ""), e.get("value"), e.get("period", ""))
                 for e in entries if e.get("value") is not None]
        if not clean:
            continue
        clean.sort(key=lambda p: (p[1] if p[1] is not None else -1e9))
        periods_present = [c[2] for c in clean if c[2]]
        all_same_period = len(set(periods_present)) <= 1
        if all_same_period:
            cat_labels = [c[0] for c in clean]
        else:
            cat_labels = [f"{c[0]} ({c[2]})" if c[2] else c[0] for c in clean]
        add_bar_chart(s, x, charts_top, col_w, charts_h, chart_title,
                      cat_labels, [c[1] for c in clean],
                      color=BRIGHT_BLUE, y_format=fmt)
    return s


def build_peer_debt(prs, data):
    return _build_peer_bars(prs, data, "peer_debt_charts",
                            "Peer comparison – Debt",
                            [("Debt / EBITDA", "debt_ebitda", "0.0\"x\""),
                             ("RCF / Net Debt", "rcf_net_debt", "0.0\"%\"")])


def build_peer_profitability(prs, data):
    return _build_peer_bars(prs, data, "peer_profitability_charts",
                            "Peer comparison – Profitability",
                            [("Revenue", "revenue", "#,##0"),
                             ("EBIT Margin", "ebit_margin", "0.0\"%\"")])


def build_peer_scatter(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_ONE_COL])
    set_text(s, 0, "Peer comparison")
    set_text(s, 18, "")
    set_text(s, 29, "")
    set_text(s, 19, "Source: Moodys.com, Company filings")
    ps = data.get("sections", {}).get("peer_scatter", {}) or {}
    bnd = ph_bounds(s, 22)
    del_ph(s, 22)
    if not bnd:
        return s
    l, t, w, h = bnd

    new_t = Inches(1.15)
    h = h + (t - new_t)
    t = new_t

    commentary = ps.get("commentary")
    if isinstance(commentary, str):
        commentary = [commentary]
    bul_h = Inches(0.6)          # was 1.0
    add_textbox_bullets(s, l, t, w, bul_h, commentary or [], size=10)
    charts_top = t + bul_h
    charts_h = h - bul_h
    col_w = int((w - Inches(0.4)) / 2)
    sc1 = ps.get("margin_vs_leverage", []) or []
    sc2 = ps.get("fcf_vs_rcf", []) or []
    _, c1 = add_scatter_chart(s, l, charts_top, col_w, charts_h,
                              "EBIT Margin vs Debt/EBITDA", sc1,
                              x_label="Debt/EBITDA",
                              y_label="EBIT Margin (%)")
    if c1:
        add_scatter_labels(s, l, charts_top, col_w, charts_h, c1)
    x2 = l + col_w + Inches(0.4)
    _, c2 = add_scatter_chart(s, x2, charts_top, col_w, charts_h,
                              "FCF/Debt vs RCF / Debt", sc2,
                              x_label="RCF / Debt", y_label="FCF / Debt")
    if c2:
        add_scatter_labels(s, x2, charts_top, col_w, charts_h, c2)
    return s


def build_scorecard(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_ONE_COL])
    set_text(s, 0, "Peer Scorecard Outcome Comparison")
    set_text(s, 18, "")
    set_text(s, 29, "")
    set_text(s, 19, "Source: Moodys.com")
    sc = data.get("sections", {}).get("scorecard", {}) or {}
    companies = sc.get("companies", []) or []
    factors = sc.get("factors", []) or []
    is_header = sc.get("is_header", []) or []
    values = sc.get("values", []) or []
    bnd = ph_bounds(s, 22)
    del_ph(s, 22)
    if not bnd or not companies or not factors:
        return s
    l, _t, w, _h = bnd
    t = Inches(1.2)
    h = Inches(5.85)
    headers = ["Factors"]
    for name in companies:
        headers.append(name)
        headers.append("")
    sub_header = [""] + ["Measure", "Score"] * len(companies)
    body = [sub_header]
    header_rows = []
    for idx, label in enumerate(factors):
        hdr = is_header[idx] if idx < len(is_header) else False
        row = [label]
        rv = values[idx] if idx < len(values) else []
        for i in range(len(companies)):
            pair = rv[i + 1] if i + 1 < len(rv) else []
            measure = pair[0] if len(pair) > 0 else ""
            score = pair[1] if len(pair) > 1 else ""
            row.append(str(measure))
            row.append(str(score))
        body.append(row)
        if hdr:
            header_rows.append(len(body) - 1)
    nrows_total = len(body) + 1
    if nrows_total <= 16:
        body_size, header_size = 9, 10
    elif nrows_total <= 19:
        body_size, header_size = 8, 9
    else:
        body_size, header_size = 7, 8
    table = add_table(s, l, t, w, h, headers, body,
                      first_col_bold=True,
                      col_widths=[2.4] + [0.65] * (len(headers) - 1),
                      header_size=header_size, body_size=body_size)
    tight_top = Emu(18000)
    tight_bot = Emu(18000)
    tight_left = Emu(36000)
    tight_right = Emu(36000)
    for r in table.rows:
        for c in r.cells:
            try:
                c.margin_top = tight_top
                c.margin_bottom = tight_bot
                c.margin_left = tight_left
                c.margin_right = tight_right
            except Exception:
                pass
    try:
        per_row = int(h / nrows_total)
        for r in table.rows:
            r.height = per_row
    except Exception:
        pass
    for row_idx in header_rows:
        for j in range(len(headers)):
            c = table.cell(row_idx, j)
            c.fill.solid()
            c.fill.fore_color.rgb = HEADER_FILL
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
    return s


def build_esg(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_TWO_COL])
    company = data.get("target_company", "")
    set_text(s, 0, "ESG Analysis")
    set_text(s, 21, "")
    set_text(s, 29, "ESG TABLE")
    set_text(s, 30, f"{company.upper()} ESG")
    set_text(s, 19, "Source: Moodys.com, Company filings")
    esg = data.get("sections", {}).get("esg_analysis", {}) or {}
    rows = esg.get("table", []) or []
    bnd_l = ph_bounds(s, 27)
    bnd_r = ph_bounds(s, 28)
    del_ph(s, 27)
    del_ph(s, 28)
    if bnd_l and rows:
        l, t, w, h = bnd_l
        headers = ["Company", "CIS", "Env.", "Social", "Gov."]
        body = [[r.get("company", ""),
                 r.get("cis", ""),
                 r.get("environmental", ""),
                 r.get("social", ""),
                 r.get("governance", "")] for r in rows]
        add_table(s, l, t, w, min(h, Inches(0.5) + Inches(0.45) * len(body)),
                  headers, body, first_col_bold=True,
                  col_widths=[1.5, 1.0, 0.9, 0.9, 0.9],
                  header_size=10, body_size=9)
    if bnd_r:
        l, t, w, h = bnd_r
        _esg_commentary = esg.get("commentary", [])
        if isinstance(_esg_commentary, str):
            _esg_commentary = [_esg_commentary] if _esg_commentary.strip() else []
        add_textbox_bullets(s, l, t, w, h, _esg_commentary, size=11)
    return s


def build_thank_you(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[L_BACK_COVER])
    set_text(s, 0, "Thank you")
    set_text(s, 94, "MOODY'S ANALYTICS")
    set_text(s, 95, data.get("sector", ""))
    set_text(s, 96, data.get("report_date", ""))
    return s


def build_disclaimer(prs):
    s = prs.slides.add_slide(prs.slide_layouts[L_DISCLAIMER])
    return s


# ---- Orchestration ---------------------------------------------------------

def build(payload: dict, out_path: Path) -> None:
    prs = Presentation(str(TEMPLATE_PATH))
    remove_all_slides(prs)

    build_cover(prs, payload)                                           # 1
    build_agenda(prs, payload)                                          # 2

    build_divider(prs, "Sector Analysis")                               # 3
    build_sector_overview(prs, payload)                                 # 4
    build_moodys_view(prs, payload)                                     # 5
    build_macro_outlook(prs, payload)                                   # 6
    build_rating_actions(prs, payload)                                  # 7

    build_divider(prs, "Company Credit Overview")                       # 8
    build_financial_analysis(prs, payload)                              # 9
    build_revenue_distribution(prs, payload)                            # 10
    build_swot(prs, payload)                                            # 11
    build_key_metrics(prs, payload)                                     # 12
    build_strategic_updates(prs, payload)                               # 13
    build_categorized(prs, payload, "news_mna",
                      "News/M&A/Leadership", "Web search")              # 14
    build_categorized(prs, payload, "external_trends",
                      "External trends / pressures / reputational risks",
                      "Web search")                                     # 15

    build_divider(prs, "Company Positioning vs. Peers")                 # 16
    build_peer_summary(prs, payload)                                    # 17
    build_peer_financials(prs, payload)                                 # 18
    build_peer_debt(prs, payload)                                       # 19
    build_peer_profitability(prs, payload)                              # 20
    build_peer_scatter(prs, payload)                                    # 21
    build_scorecard(prs, payload)                                       # 22
    build_esg(prs, payload)                                             # 23

    build_thank_you(prs, payload)                                       # 24
    build_disclaimer(prs)                                               # 25

    for slide in prs.slides:
        clear_empty_placeholders(slide)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> int:
    if len(sys.argv) != 3:
        print("Usage: build_pptx.py <payload.json> <output.pptx>",
              file=sys.stderr)
        return 2
    with Path(sys.argv[1]).open("r", encoding="utf-8") as f:
        payload = json.load(f)
    build(payload, Path(sys.argv[2]))
    print(f"Wrote {sys.argv[2]}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
